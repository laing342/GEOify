#!/usr/bin/env python3
"""
Generate a professional PPT presentation for the 轻遇星 GEO internship assignment.
Based on the GEOify fork project workflow.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette — purple theme matching the demo site
PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0xB8, 0x94)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    """Add a solid background color to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, spacing=Pt(8)):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
        p.level = 0
    return tf


# ============================================================
# Slide 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Decorative bar
add_shape_bg(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             '轻遇星 GEO 优化实习作业', font_size=44, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
             'Generative Engine Optimization · 6 引擎基线测评 + Demo 站实战',
             font_size=22, color=RGBColor(0xAA, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             '基于 Geoify 魔改 · 轻遇星品牌 · 1 周实战',
             font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             '2026-07-24', font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: Project Background
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '项目背景：轻遇星 & GEO', font_size=32, bold=True)

# Left panel
add_shape_bg(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8), LIGHT_GRAY, 0.05)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(0.6),
             '🏢  轻遇星 (iCityMuze)', font_size=22, bold=True, color=DARK_BG)
add_bullet_slide(slide, Inches(1.2), Inches(2.9), Inches(5), Inches(3.5), [
    'AI 驱动的 O2O 社交陪护与灵活就业平台',
    '核心服务：城市伴游、商务接待、景点讲解、',
    '　　　　　情绪陪伴、本地轻社交小聚',
    '三大 AI 能力：达人数字分身、',
    '　　　　　　灵魂契合推荐、AI 获客助理',
    '重点城市：南京、徐州、苏州',
    '达人侧：普通人可申请，灵活就业',
], font_size=16, color=RGBColor(0x33, 0x33, 0x44))

# Right panel
add_shape_bg(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.8), LIGHT_GRAY, 0.05)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.6),
             '🔍  什么是 GEO？', font_size=22, bold=True, color=DARK_BG)
add_bullet_slide(slide, Inches(7.3), Inches(2.9), Inches(5), Inches(3.5), [
    'GEO = Generative Engine Optimization',
    '　　 生成式引擎优化',
    '',
    'SEO 追求: 搜索引擎排名靠前',
    'GEO 追求: 被 AI 回答引用',
    '',
    'AI 引用路径:',
    '　能被搜到 → 排序 → 可信度加权 → 引用',
    '',
    '中国特有: 豆包/DeepSeek/Kimi/元宝/',
    '　　　　　百度AI/智谱 需各自索引收录',
], font_size=16, color=RGBColor(0x33, 0x33, 0x44))

# ============================================================
# Slide 3: Tech Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '技术架构：Geoify 魔改方案', font_size=32, bold=True)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
             '不动原有代码 · 新增 4 个 TypeScript 源文件 · 扩展为静态站生成器',
             font_size=18, color=MED_GRAY)

# File boxes
files_data = [
    ('config.ts', '站点配置文件\n品牌·FAQ·城市·信任体系', GREEN),
    ('jsonld-injector.ts', 'JSON-LD 注入器\nFAQPage + Organization', RGBColor(0x00, 0x96, 0xC7)),
    ('llms-generator.ts', 'llms.txt 生成器\nAI 引擎友好站点导航', ORANGE),
    ('generate-site.ts', '静态站生成器\nHTML + robots.txt + sitemap', PURPLE),
]

for i, (name, desc, color) in enumerate(files_data):
    left = Inches(1 + i * 3)
    # Card bg
    add_shape_bg(slide, left, Inches(2.8), Inches(2.6), Inches(2.8), LIGHT_GRAY, 0.05)
    # Color bar on top
    add_shape_bg(slide, left, Inches(2.8), Inches(2.6), Inches(0.08), color)
    # File name
    add_text_box(slide, Inches(0.2) + left, Inches(3.1), Inches(2.2), Inches(0.6),
                 name, font_size=20, bold=True, color=DARK_BG)
    # Description
    add_text_box(slide, Inches(0.2) + left, Inches(3.7), Inches(2.2), Inches(1.8),
                 desc, font_size=15, color=RGBColor(0x55, 0x55, 0x66))

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
             'CLI 入口: src/cli.ts → 注册 generate-site 命令 → 一键生成完整 Demo 站',
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 4: Data Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '数据流转：从配置到部署', font_size=32, bold=True)

# Flow steps
steps = [
    ('1', 'config.ts\n配置文件', GREEN),
    ('2', 'generate-site.ts\n站点生成器', PURPLE),
    ('3', 'HTML 页面\n×6 个页面', RGBColor(0x00, 0x96, 0xC7)),
    ('4', 'JSON-LD\n结构化数据', ORANGE),
    ('5', 'GitHub Pages\n一键部署', GREEN),
]

for i, (num, label, color) in enumerate(steps):
    left = Inches(0.8 + i * 2.5)
    top = Inches(2.5)
    # Circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.5), top, Inches(1.2), Inches(1.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(10)

    # Label
    add_text_box(slide, left, top + Inches(1.4), Inches(2.2), Inches(1),
                 label, font_size=14, color=DARK_BG, alignment=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            left + Inches(1.85), top + Inches(0.3),
            Inches(0.5), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xEE)
        arrow.line.fill.background()

# Output summary
add_shape_bg(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(1.8), LIGHT_GRAY, 0.05)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
             '📦  单次命令输出:', font_size=18, bold=True, color=DARK_BG)
add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.8),
             '6 个 HTML 页面 ＋ robots.txt ＋ sitemap.xml ＋ llms.txt ＋ FAQPage JSON-LD ＋ Organization JSON-LD',
             font_size=16, color=RGBColor(0x55, 0x55, 0x66))

# ============================================================
# Slide 5: Demo Site Pages
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             'Demo 站点：6 个页面 + 3 个系统文件', font_size=32, bold=True)

# Page cards
pages = [
    ('🏠', 'index.html', '品牌定位页\n三大 AI 能力\nFAQ 预览'),
    ('❓', 'faq.html', '6 条 FAQ\n答案先行 ≤80 字\nFAQPage JSON-LD'),
    ('🛡️', 'trust.html', '五层安全防线\n红线制度\n阳光陪伴声明'),
    ('🏙️', 'cities/', '南京·徐州·苏州\n价格对比表\n地域长尾覆盖'),
    ('🤖', 'llms.txt', 'AI 引擎导航\n站点结构描述\n标准格式'),
    ('🗺️', '系统文件', 'robots.txt\nsitemap.xml\n6 爬虫放行规则'),
]

for i, (icon, name, desc) in enumerate(pages):
    left = Inches(0.5 + (i % 3) * 4.2)
    top = Inches(1.8 + (i // 3) * 2.7)
    add_shape_bg(slide, left, top, Inches(3.8), Inches(2.3), LIGHT_GRAY, 0.05)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(3.2), Inches(0.5),
                 f'{icon}  {name}', font_size=20, bold=True, color=PURPLE)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.2), Inches(1.4),
                 desc, font_size=15, color=RGBColor(0x55, 0x55, 0x66))

# ============================================================
# Slide 6: JSON-LD Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             'JSON-LD 结构化数据：FAQPage Schema', font_size=32, bold=True)

# Code block
add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.2),
             RGBColor(0x2D, 0x2D, 0x3F), 0.03)

code_lines = [
    '{',
    '  "@context": "https://schema.org",',
    '  "@type": "FAQPage",',
    '  "mainEntity": [',
    '    {',
    '      "@type": "Question",',
    '      "name": "轻遇星是做什么的？",',
    '      "acceptedAnswer": {',
    '        "@type": "Answer",',
    '        "text": "轻遇星（iCityMuze）是AI驱动的..."',
    '      }',
    '    },',
    '    // ... 共 6 条 Q&A ...',
    '  ],',
    '  "about": {',
    '    "@type": "Organization",',
    '    "name": "轻遇星"',
    '  }',
    '}',
]
for i, line in enumerate(code_lines):
    color = WHITE if i < 2 or i >= len(code_lines) - 3 else RGBColor(0xAA, 0xCC, 0xFF)
    add_text_box(slide, Inches(1.2), Inches(2.0 + i * 0.28), Inches(6.8), Inches(0.3),
                 line, font_size=13, color=color, font_name='Consolas')

# Right panel: checklist
add_shape_bg(slide, Inches(8.8), Inches(1.8), Inches(3.8), Inches(5.2), PURPLE, 0.05)
add_text_box(slide, Inches(9.2), Inches(2.0), Inches(3.2), Inches(0.6),
             '✅ 自检通过', font_size=22, bold=True)
check_items = [
    '✓ FAQPage 类型正确',
    '✓ 6 条 Question+Answer',
    '✓ @context schema.org',
    '✓ Organization 关联',
    '✓ 标记文字 = 可见文字',
    '✓ Schema 验证器通过',
    '✓ 所有字段合法',
    '✓ 无必填字段缺失',
]
for i, item in enumerate(check_items):
    add_text_box(slide, Inches(9.2), Inches(2.7 + i * 0.45), Inches(3.2), Inches(0.4),
                 item, font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF))

# ============================================================
# Slide 7: Baseline Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '任务 1：6 引擎基线测评结果', font_size=32, bold=True)

add_text_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.6),
             '测评时间：2026-07-24 · 40 题 × 6 引擎 = 240 次查询 · 站点上线 < 24h',
             font_size=16, color=MED_GRAY)

# Table header
table_data = [
    ['引擎', '品牌提及率', '被引用率', '正面率', '准确率', '首位提及率'],
    ['豆包', '0%', '0%', 'N/A', 'N/A', 'N/A'],
    ['Kimi', '0%', '0%', 'N/A', 'N/A', 'N/A'],
    ['DeepSeek', '0%', '0%', 'N/A', 'N/A', 'N/A'],
    ['腾讯元宝', '0%', '0%', 'N/A', 'N/A', 'N/A'],
    ['百度 AI', '0%', '0%', 'N/A', 'N/A', 'N/A'],
    ['智谱清言', '0%', '0%', 'N/A', 'N/A', 'N/A'],
]

table = slide.shapes.add_table(
    len(table_data), len(table_data[0]),
    Inches(1), Inches(2.4), Inches(11.3), Inches(3.2)
).table

for r, row in enumerate(table_data):
    for c, cell_text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.name = 'Microsoft YaHei'
            paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_BG
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PURPLE
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

# Conclusion
add_shape_bg(slide, Inches(1), Inches(5.9), Inches(11.3), Inches(1.2), LIGHT_GRAY, 0.05)
add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10.5), Inches(1),
             '💡  结论：轻遇星当前 AI 可见度为 0 —— 这本身就是正确的基线结果。\n'
             '　　　站点刚上线、未索引、无外链。Week 2 复测预期开始出现增长。',
             font_size=16, color=RGBColor(0x33, 0x33, 0x44))

# ============================================================
# Slide 8: SEO Indexing Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '收录加速：搜索引擎优化策略', font_size=32, bold=True)

strategies = [
    ('IndexNow API 即时推送', 'HTTP 202 Accepted\nBing + Yandex 实时收录\n无需站长平台注册', GREEN),
    ('robots.txt 放行规则', 'Baiduspider / YisouSpider\nBytespider / Bingbot\nGPTBot / PerplexityBot / CCBot', PURPLE),
    ('sitemap.xml', '6 个 URL 全部列出\n每周更新频率\n标准 XML 格式', ORANGE),
    ('llms.txt', 'AI 引擎专用站点地图\n符合 llmstxt.org 标准\n国产 AI 引擎兼容', RGBColor(0x00, 0x96, 0xC7)),
]

for i, (title, desc, color) in enumerate(strategies):
    left = Inches(0.8 + i * 3.15)
    add_shape_bg(slide, left, Inches(2.2), Inches(2.8), Inches(4.2), LIGHT_GRAY, 0.05)
    add_shape_bg(slide, left, Inches(2.2), Inches(2.8), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.8),
                 title, font_size=18, bold=True, color=DARK_BG)
    add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(2.4), Inches(2.8),
                 desc, font_size=14, color=RGBColor(0x55, 0x55, 0x66))

# ============================================================
# Slide 9: Content Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '任务 2：内容策略与信息架构', font_size=32, bold=True)

# 5 columns
columns = [
    ('品牌/定位页', 'index.html', GREEN,
     ['轻遇星是做什么的', '轻遇星靠谱吗', '轻遇星怎么样']),
    ('FAQ 页', 'faq.html', PURPLE,
     ['怎么成为达人', '做地陪赚多少钱', '轻遇星有哪些服务']),
    ('城市指南页', 'cities/*.html', RGBColor(0x00, 0x96, 0xC7),
     ['南京本地向导', '苏州商务接待', '徐州伏羊节攻略']),
    ('信任/安全页', 'trust.html', ORANGE,
     ['轻遇星安全吗', '五层防线', '怎么保障隐私']),
    ('达人入驻页', 'join.html', MED_GRAY,
     ['入驻条件', '零押金/免培训', '审核 1-3 天']),
]

for i, (title, url, color, queries) in enumerate(columns):
    left = Inches(0.5 + i * 2.55)
    add_shape_bg(slide, left, Inches(2.0), Inches(2.3), Inches(4.5), LIGHT_GRAY, 0.05)
    add_shape_bg(slide, left, Inches(2.0), Inches(2.3), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.15), Inches(2.2), Inches(2), Inches(0.5),
                 title, font_size=17, bold=True, color=DARK_BG)
    add_text_box(slide, left + Inches(0.15), Inches(2.55), Inches(2), Inches(0.4),
                 url, font_size=12, color=color)
    for j, q in enumerate(queries):
        add_text_box(slide, left + Inches(0.15), Inches(3.05 + j * 0.5), Inches(2), Inches(0.45),
                     f'🔍 {q}', font_size=13, color=RGBColor(0x33, 0x33, 0x44))

add_shape_bg(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), LIGHT_GRAY, 0.05)
add_text_box(slide, Inches(1), Inches(6.75), Inches(11), Inches(0.4),
             '💡 每个页面直接回答用户问 AI 的原话 · 答案先行 ≤80 字 · 含地域长尾关键词 · 每 100 字 1 个可核实数据',
             font_size=15, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 10: Metrics & Iteration
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '任务 5：度量、迭代与 1 个月实战', font_size=32, bold=True)

# Left: KPI
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.5),
             '📊  5 项核心指标 (K1-K5)', font_size=20, bold=True, color=DARK_BG)

kpis = [
    'K1 品牌提及率:  ≥ 80% (5/6 引擎出现品牌名)',
    'K2 被引用率 (SoM):  ≥ 30% (40 题中被引用)',
    'K3 正面率:  ≥ 90% (非负面引用占比)',
    'K4 准确率:  ≥ 95% (事实与实际一致)',
    'K5 首位提及率:  ≥ 20% (排第一的占比)',
]
add_bullet_slide(slide, Inches(1), Inches(2.3), Inches(5.5), Inches(2.5),
                 kpis, font_size=15, color=RGBColor(0x33, 0x33, 0x44))

# Right: weekly rhythm
add_text_box(slide, Inches(7), Inches(1.7), Inches(5), Inches(0.5),
             '📅  周度迭代节奏', font_size=20, bold=True, color=DARK_BG)
rhythm = [
    '周一: 40 题 × 6 引擎复测',
    '周三: 分析引用变化 + 对比上周',
    '周五: 修改配置 → 重新生成站点',
    '周日: 周报 1 页',
]
add_bullet_slide(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(2.5),
                 rhythm, font_size=15, color=RGBColor(0x33, 0x33, 0x44))

# Bottom: 1 month real world
add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.2), LIGHT_GRAY, 0.05)
add_text_box(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.5),
             '🚀  1 个月 + 真实网站推进路线 (加分思考题)', font_size=18, bold=True, color=DARK_BG)
month_items = [
    'Week 1: robots.txt + sitemap + 站长工具 + 部署 GEO 结构（FAQ+JSON-LD+城市页+信任页）',
    'Week 2: IndexNow 推送 + 外链 + 初始内容库填充',
    'Week 3: ≥5 篇城市指南 + 对比表 + 开始追踪引用',
    'Week 4: A/B 测试 + 竞品对比页 + 根据数据迭代',
    '⚠️  坑: github.io 权重低 → 绑定自有域名 → 做外链 · 收录 ≠ 引用，间隔可能 1-4 周',
]
add_bullet_slide(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(1.5),
                 month_items, font_size=13, color=RGBColor(0x55, 0x55, 0x66), spacing=Pt(4))

# ============================================================
# Slide 11: Deliverables Checklist
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), PURPLE)
add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(1),
             '交付物清单 & 成果展示', font_size=32, bold=True)

deliverables = [
    ('📄', '任务 1', '基线测评报告', '6 引擎 × 40 题'),
    ('📐', '任务 2', '内容策略与信息架构', '页面地图 + 关键词表'),
    ('✍️', '任务 3', '样例内容 3 篇', 'FAQ + 城市指南 + 信任页'),
    ('🌐', '任务 4', 'Demo 站 + 自检', 'laing342.github.io/GEOify'),
    ('📊', '任务 5', '度量与迭代计划', 'KPI + 周迭代 + 1 月路线'),
    ('📋', '任务 6', 'README 说明', '文件索引 + 思路 + 用时'),
]

for i, (icon, task, name, desc) in enumerate(deliverables):
    y = Inches(1.8 + i * 0.85)
    add_shape_bg(slide, Inches(1.5), y, Inches(10.3), Inches(0.7), PURPLE, 0.03)
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(0.6), Inches(0.5),
                 icon, font_size=20, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.1), Inches(2), Inches(0.5),
                 f'{task}', font_size=16, bold=True, color=RGBColor(0xCC, 0xCC, 0xFF))
    add_text_box(slide, Inches(4.5), y + Inches(0.1), Inches(3), Inches(0.5),
                 name, font_size=16, color=WHITE)
    add_text_box(slide, Inches(7.8), y + Inches(0.1), Inches(3.5), Inches(0.5),
                 desc, font_size=14, color=RGBColor(0xAA, 0xAA, 0xBB))

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             '✅ 全部 6 项交付物已完成 · 合规红线自检通过 · 零改动原有 Geoify 代码',
             font_size=15, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 12: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_bg(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             '谢谢！', font_size=48, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             '轻遇星 GEO 优化实习作业', font_size=22, color=RGBColor(0xAA, 0xAA, 0xBB),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             'Demo 站: https://laing342.github.io/GEOify/\n'
             '源码:  https://github.com/laing342/GEOify',
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             '基于 Geoify 魔改 · 2026-07-24', font_size=14, color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'docs', '轻遇星GEO优化实习作业.pptx')
prs.save(output_path)
print(f'PPT saved to: {output_path}')
